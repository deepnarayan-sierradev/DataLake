"""
Enterprise Data Lake Platform — Presentation Generator
Theme: exactly matches AI_Weekly_Report_01152026.pptx
  - Pure white background on every slide (NO dark slides)
  - Right-side vertical gradient strip: orange(top) → pale-lime(bottom)
  - Title: ~38pt bold, colour #383E48 (charcoal), font Open Sauce Bold
  - Body: 13-15pt, colour #222222, font Open Sauce
  - Primary accent colour: #156082 (deep blue)
  - Secondary: #E97132 (orange), #0F9ED5 (sky blue), #196B24 (green)
  - NO overlapping text — all boxes sized generously

Run: python scripts/generate_presentation.py
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import tempfile
from pathlib import Path

CHARCOAL  = RGBColor(0x38, 0x3E, 0x48)  # title text  (#383E48)
DEEP_BLUE = RGBColor(0x15, 0x60, 0x82)  # accent1
ORANGE    = RGBColor(0xE9, 0x71, 0x32)  # accent2
SKY       = RGBColor(0x0F, 0x9E, 0xD5)  # accent4
GREEN     = RGBColor(0x19, 0x6B, 0x24)  # accent3
PURPLE    = RGBColor(0x6A, 0x3F, 0xC0)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF8, 0xF8, 0xF8)
RULE_C    = RGBColor(0xD0, 0xD0, 0xD0)
TEXT      = RGBColor(0x22, 0x22, 0x22)
MUTED     = RGBColor(0x66, 0x66, 0x66)
GREEN_OK  = RGBColor(0x19, 0x6B, 0x24)
AMBER     = RGBColor(0xFF, 0xA5, 0x00)
RED       = RGBColor(0xC0, 0x20, 0x20)
LITE_BLUE = RGBColor(0xE3, 0xEE, 0xF8)
LITE_GRNN = RGBColor(0xE8, 0xF5, 0xE9)

FH = "Open Sauce Bold"
FB = "Open Sauce"
FM = "Courier New"

SW = Inches(13.33)
SH = Inches(7.5)

GX = Inches(10.9)   # strip left edge
GW = Inches(2.43)   # strip width
_STRIP = str(Path(tempfile.gettempdir()) / "gradient_strip.png")

CX  = Inches(0.55)   # left margin
CW  = Inches(10.1)   # content width (leaves gap before strip)
TY  = Inches(0.42)   # title top
TH  = Inches(0.82)   # title height
RY  = Inches(1.30)   # rule y
STY = Inches(1.38)   # subtitle top
STH = Inches(0.35)   # subtitle height
BY  = Inches(1.82)   # body start y
BH  = Inches(5.15)   # body height (to y=6.97)
FTY = Inches(7.10)   # footer y


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _white(s):
    f = s.background.fill; f.solid(); f.fore_color.rgb = WHITE

def _strip(s):
    if os.path.exists(_STRIP):
        s.shapes.add_picture(_STRIP, GX, Inches(0), GW, SH)

def _box(s, l, t, w, h, fill=None, lc=None, lw=Pt(0)):
    shp = s.shapes.add_shape(1, l, t, w, h)
    shp.line.width = lw
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if lc:
        shp.line.color.rgb = lc
    else:
        shp.line.fill.background()
    return shp

def _t(s, text, l, t, w, h, sz=Pt(14), bold=False, color=TEXT,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    txb = s.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = sz; r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font or (FH if bold else FB)
    return txb

def _blt(s, items, l, t, w, h, sz=Pt(13), color=TEXT, bchar="•  ",
         hdg=None, hc=None, hsz=None):
    """Multi-line bullet/list text box. hdg = optional bold heading line."""
    txb = s.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame; tf.word_wrap = True
    first = True
    if hdg:
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = hdg
        r.font.size = hsz or (sz + Pt(2))
        r.font.bold = True
        r.font.color.rgb = hc or DEEP_BLUE
        r.font.name = FH
        first = False
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]; first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = f"{bchar}{item}"
        r.font.size = sz; r.font.color.rgb = color; r.font.name = FB
    return txb

def _rule(s, y, color=RULE_C, thick=Pt(1.5)):
    _box(s, CX, y, GX - CX - Inches(0.15), thick, fill=color)

def _vline(s, x, y_top, height, color=RULE_C, w=Pt(1)):
    _box(s, x, y_top, w, height, fill=color)

def _header(s, title, subtitle=None):
    """Standard white-background header (title + rule + optional subtitle)."""
    _t(s, title, CX, TY, CW, TH, sz=Pt(36), bold=True, color=CHARCOAL, font=FH)
    _rule(s, RY, color=DEEP_BLUE, thick=Pt(2))
    if subtitle:
        _t(s, subtitle, CX, STY, CW, STH, sz=Pt(14), color=DEEP_BLUE, font=FB)

def _footer(s, txt="Enterprise Data Lake Platform  ·  July 2026"):
    _t(s, txt, CX, FTY, CW, Inches(0.32), sz=Pt(10), color=MUTED, font=FB)


def s01_title(prs):
    s = _blank(prs); _white(s); _strip(s)

    _t(s, "Enterprise Data Lake Platform",
       CX, Inches(0.65), Inches(10.0), Inches(1.2),
       sz=Pt(44), bold=True, color=CHARCOAL, font=FH)

    _t(s, "Automated  ·  Metadata-Driven  ·  Cloud-Native  ·  Governed",
       CX, Inches(1.95), Inches(10.0), Inches(0.42),
       sz=Pt(16), color=DEEP_BLUE, font=FB)

    _rule(s, Inches(2.45), color=DEEP_BLUE, thick=Pt(2))

    _t(s, "A unified platform connecting Salesforce CRM, MySQL RDS, and Sage ERP\n"
          "into a single trusted, analytics-ready data foundation — built on AWS.",
       CX, Inches(2.58), Inches(9.8), Inches(0.9),
       sz=Pt(16), color=TEXT, font=FB)

    badges = [(GREEN_OK, "Dev  ✅  Live"), (AMBER, "Staging  🔲  Next"), (RED, "Production  🔲  Planned")]
    for i, (c, lbl) in enumerate(badges):
        bx = CX + i * Inches(3.0)
        _box(s, bx, Inches(3.75), Inches(2.75), Inches(0.46), fill=c)
        _t(s, lbl, bx + Inches(0.1), Inches(3.79), Inches(2.55), Inches(0.38),
           sz=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _t(s, "Platform Engineering  ·  July 2026",
       CX, FTY, Inches(7.0), Inches(0.32), sz=Pt(11), color=MUTED, font=FB)


def s02_problem(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "The Problem We Solved",
            "Data lived in silos — disconnected, delayed, and insecure")

    cx  = [CX, Inches(3.15), Inches(5.8),  Inches(8.3)]
    cw  = [Inches(2.45), Inches(2.45), Inches(2.35), Inches(2.3)]
    hdrs = ["System", "What it held", "Who used it", "Pain point"]
    ht = BY
    for h, x, w in zip(hdrs, cx, cw):
        _box(s, x, ht, w, Inches(0.4), fill=DEEP_BLUE)
        _t(s, h, x+Inches(0.08), ht+Inches(0.07), w-Inches(0.12), Inches(0.3),
           sz=Pt(12), bold=True, color=WHITE, font=FH)

    rows = [
        ("Salesforce CRM",  "Accounts & contacts",          "Sales, CS",        "Manual exports; no audit trail"),
        ("MySQL RDS",       "Contracts & orders",            "Operations",       "24–72 hr delay; brittle scripts"),
        ("Sage Intacct/X3", "AR/AP invoices, vendors",       "Finance",          "Disconnected from CRM"),
        ("All systems",     "Credentials stored in scripts", "Engineering",      "Security risk, no rotation"),
        ("All systems",     "Same customer = 3 records",     "Analytics",        "No single source of truth"),
    ]
    for ri, (row, bg) in enumerate(zip(rows, [WHITE, LIGHT_BG] * 10)):
        rt = BY + Inches(0.4) + ri * Inches(0.5)
        for val, x, w in zip(row, cx, cw):
            _box(s, x, rt, w, Inches(0.48), fill=bg)
            _t(s, val, x+Inches(0.08), rt+Inches(0.07), w-Inches(0.12), Inches(0.38),
               sz=Pt(11.5), color=TEXT, font=FB)

    _box(s, CX, Inches(5.0), Inches(10.1), Inches(0.5), fill=DEEP_BLUE)
    _t(s, "Result: stale data  ·  no audit trail  ·  no single source of truth  ·  weeks to onboard any new data source",
       CX+Inches(0.15), Inches(5.07), Inches(9.8), Inches(0.38),
       sz=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FH)

    pairs = [("24–72 hr delay","1–4 hr automated"),
             ("3 disconnected views","1 golden record"),
             ("No audit trail","Full lineage"),
             ("Credentials in scripts","AWS Secrets Manager")]
    for i, (bef, aft) in enumerate(pairs):
        bx = CX + i * Inches(2.5)
        _box(s, bx, Inches(5.68), Inches(1.15), Inches(0.65), fill=RGBColor(0xFF,0xEE,0xEE))
        _t(s, bef, bx+Inches(0.05), Inches(5.72), Inches(1.05), Inches(0.58),
           sz=Pt(10), color=RED, align=PP_ALIGN.CENTER, font=FB)
        _t(s, "→", bx+Inches(1.18), Inches(5.88), Inches(0.22), Inches(0.3),
           sz=Pt(12), bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
        _box(s, bx+Inches(1.43), Inches(5.68), Inches(1.0), Inches(0.65), fill=RGBColor(0xE8,0xF5,0xE9))
        _t(s, aft, bx+Inches(1.48), Inches(5.72), Inches(0.9), Inches(0.58),
           sz=Pt(10), color=GREEN_OK, align=PP_ALIGN.CENTER, font=FB)

    _footer(s)


def s03_overview(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "What We Built",
            "A fully automated, metadata-driven data lake on AWS")

    cards = [
        (DEEP_BLUE, "Multi-Source\nConnectors",
         "Salesforce  ·  MySQL RDS\nSage Intacct  ·  Sage X3\nOAuth 2.0 / REST / JWT"),
        (SKY, "4-Stage\nPipeline",
         "Extract  →  Transform\n→  Resolve  →  Publish\nAWS Step Functions"),
        (PURPLE, "Entity\nResolution",
         "One golden record per\ncustomer across all\nsource systems"),
        (ORANGE, "Governed\nAnalytics",
         "SQL via Athena\nPII masking  ·  Lineage\nQuality report per run"),
    ]
    cw = Inches(2.4)
    for (c, title, body), x in zip(cards, [CX, CX+Inches(2.55), CX+Inches(5.1), CX+Inches(7.65)]):
        _box(s, x, BY, cw, Inches(4.3), fill=c)
        _t(s, title, x+Inches(0.15), BY+Inches(0.15), cw-Inches(0.25), Inches(0.85),
           sz=Pt(17), bold=True, color=WHITE, font=FH)
        _box(s, x+Inches(0.08), BY+Inches(1.05), cw-Inches(0.16), Pt(1), fill=WHITE)
        _t(s, body, x+Inches(0.15), BY+Inches(1.18), cw-Inches(0.25), Inches(3.0),
           sz=Pt(14), color=WHITE, font=FB)

    _box(s, CX, Inches(6.5), Inches(10.1), Inches(0.68), fill=DEEP_BLUE)
    stats = [("4 sources","connected"), ("35,971+","records live"), ("0 code","to add a source"),
             ("< 4 hrs","end-to-end"), ("~$469/mo","AWS cost")]
    sw = Inches(10.1) / len(stats)
    for i, (num, lbl) in enumerate(stats):
        sx = CX + i * sw
        _t(s, num, sx+Inches(0.05), Inches(6.53), sw-Inches(0.08), Inches(0.35),
           sz=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FH)
        _t(s, lbl, sx+Inches(0.05), Inches(6.87), sw-Inches(0.08), Inches(0.28),
           sz=Pt(11), color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER, font=FB)

    _footer(s)


def s04_architecture(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "Architecture Overview",
            "Event-driven  ·  Serverless  ·  Three governed data layers")

    layers = [
        (CX, BY, Inches(10.1), Inches(0.72),
         LITE_BLUE, DEEP_BLUE,
         "SOURCE SYSTEMS",
         "Salesforce CRM  ·  MySQL RDS  ·  Sage Intacct  ·  Sage X3  ·  (NetSuite — roadmap)"),

        (CX, BY+Inches(0.78), Inches(10.1), Inches(0.72),
         RGBColor(0xF3,0xEC,0xFF), PURPLE,
         "ORCHESTRATION",
         "Amazon EventBridge Scheduler  →  AWS Step Functions  ·  DynamoDB (config, watermark, audit)"),
    ]
    for lx, ly, lw, lh, lf, lb_c, lt, ld in layers:
        _box(s, lx, ly, lw, lh, fill=lf, lc=lb_c, lw=Pt(1.5))
        _t(s, lt, lx+Inches(0.12), ly+Inches(0.07), lw-Inches(0.18), Inches(0.28),
           sz=Pt(10), bold=True, color=lb_c, font=FH)
        _t(s, ld, lx+Inches(0.12), ly+Inches(0.37), lw-Inches(0.18), Inches(0.32),
           sz=Pt(11.5), color=TEXT, font=FB)

    dlayers = [
        (CX,                     "RAW LAYER  (S3)",    LITE_GRNN,           GREEN,
         "Immutable · Parquet\nObject Lock\nSchema snapshots\nDrift detection"),
        (CX+Inches(3.45),        "CURATED LAYER  (S3)", RGBColor(0xEE,0xEE,0xFF), PURPLE,
         "Field-mapped · Masked\nSCD Type 1 merge\nCanonical names\nGlue-catalogued"),
        (CX+Inches(6.9),         "ANALYTICS LAYER  (S3 + Glue)", RGBColor(0xFF,0xF3,0xE8), ORANGE,
         "Golden records\nAthena-queryable\nPartitioned by date\nLake Formation"),
    ]
    dlw = Inches(3.15)
    for dx, dt, df, dc, dd in dlayers:
        _box(s, dx, BY+Inches(1.56), dlw, Inches(2.55), fill=df, lc=dc, lw=Pt(1.5))
        _t(s, dt, dx+Inches(0.12), BY+Inches(1.62), dlw-Inches(0.18), Inches(0.3),
           sz=Pt(10), bold=True, color=dc, font=FH)
        _t(s, dd, dx+Inches(0.12), BY+Inches(1.97), dlw-Inches(0.18), Inches(2.1),
           sz=Pt(12), color=TEXT, font=FB)

    for ax in [CX+Inches(3.25), CX+Inches(6.7)]:
        _t(s, "→", ax, BY+Inches(2.5), Inches(0.32), Inches(0.4),
           sz=Pt(22), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    lambdas = [(CX, "Lambda 1\nExtraction", DEEP_BLUE),
               (CX+Inches(2.6), "Lambda 2\nTransformation", PURPLE),
               (CX+Inches(5.2), "Lambda 3\nEntity Resolution", GREEN),
               (CX+Inches(7.8), "Lambda 4\nAnalytics Publisher", ORANGE)]
    lbw = Inches(2.2)
    for lx, lt, lc in lambdas:
        _box(s, lx, BY+Inches(4.2), lbw, Inches(0.65), fill=lc)
        _t(s, lt, lx+Inches(0.08), BY+Inches(4.24), lbw-Inches(0.12), Inches(0.58),
           sz=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FH)

    _footer(s)


def s05_sources(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "Data Source Connectors",
            "Pluggable adapter pattern — add a new source without changing the pipeline")

    srcs = [
        (DEEP_BLUE, "✅", "Salesforce CRM",
         "OAuth 2.0 client_credentials\nBulk API 2.0 (async CSV)\nEntities: Account, Contact\nWatermark: SystemModstamp"),
        (GREEN,     "✅", "MySQL RDS",
         "AWS Secrets Manager creds\nVPC + NAT Gateway\nEntity: Contracts (incremental)\nWatermark: ModifiedOn"),
        (PURPLE,    "✅", "Sage Intacct",
         "REST API, OAuth 2.0\nEntities: Customer, Vendor,\nAR Invoice, AP Bill\nWatermark: auditInfo.modifiedAt"),
        (ORANGE,    "✅", "Sage X3",
         "REST API, OAuth 2.0\nEntities: Customer, Supplier\nSame SageConnector adapter\nWatermark incremental"),
        (MUTED,     "🔲", "NetSuite ERP",
         "OAuth 1.0a token-based\nFuture — config only\nNo code change needed\nCredentials pre-reserved"),
    ]
    cw = Inches(1.95)
    for i, (c, st, nm, body) in enumerate(srcs):
        sx = CX + i * Inches(2.06)
        active = st == "✅"
        _box(s, sx, BY, cw, Inches(5.1), fill=c if active else LIGHT_BG,
             lc=None if active else RULE_C, lw=Pt(1))
        _t(s, st, sx+Inches(0.1), BY+Inches(0.12), cw-Inches(0.15), Inches(0.36),
           sz=Pt(18), color=WHITE if active else MUTED)
        _t(s, nm, sx+Inches(0.1), BY+Inches(0.52), cw-Inches(0.15), Inches(0.48),
           sz=Pt(13), bold=True, color=WHITE if active else CHARCOAL, font=FH)
        _t(s, body, sx+Inches(0.1), BY+Inches(1.05), cw-Inches(0.15), Inches(3.9),
           sz=Pt(11), color=RGBColor(0xEE,0xEE,0xEE) if active else MUTED, font=FB)

    _box(s, CX, Inches(7.02), Inches(10.1), Inches(0.38), fill=DEEP_BLUE)
    _t(s, "All connectors implement ExtractionConnector  ·  3 methods to register a new source",
       CX+Inches(0.15), Inches(7.08), Inches(9.8), Inches(0.28),
       sz=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FH)


def s06_extraction(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "Pipeline Stage 1 — Extraction",
            "Source systems  →  S3 Raw Layer")

    C1X, C1W = CX, Inches(4.75)
    C2X, C2W = CX+Inches(5.2), Inches(4.75)
    _vline(s, CX+Inches(5.0), BY, Inches(3.3), color=RULE_C)

    _blt(s, ["Config from DynamoDB: entity_id, load_type, watermark_field",
             "Credentials from Secrets Manager — never in code",
             "Metadata discovery validates schema before querying",
             "Query: full load  or  incremental (watermark window)",
             "Records streamed one-at-a-time — no large memory use",
             "Schema snapshot saved to S3 after each extraction",
             "Drift detection: added / removed / type-changed columns",],
         C1X, BY, C1W, Inches(3.2),
         sz=Pt(12.5), color=TEXT, hdg="Processing Steps", hc=DEEP_BLUE)

    _blt(s, ["Circuit breaker (DynamoDB-backed) — stops broken sources",
             "Exponential backoff + jitter on transient errors",
             "Dead-letter queue (SQS) for replay on failures",
             "First-run incremental: epoch (1970) lower bound",
             "S3 Object Lock on raw layer — immutable audit trail",
             "IAM least-privilege scoped to exact S3 prefixes",
             "Run audit log entry written to DynamoDB per run",],
         C2X, BY, C2W, Inches(3.2),
         sz=Pt(12.5), color=TEXT, hdg="Reliability & Security", hc=ORANGE)

    _box(s, CX, Inches(5.3), Inches(10.1), Inches(0.56), fill=RGBColor(0xF0,0xF4,0xF8))
    _t(s, "S3 path:  s3://edl-raw-087972550871/raw/{source_id}/{entity_id}/extraction_date=YYYY-MM-DD/part-NNNNN.parquet",
       CX+Inches(0.15), Inches(5.36), Inches(9.8), Inches(0.48),
       sz=Pt(11), color=DEEP_BLUE, font=FM)

    _footer(s)


def s07_transformation(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "Pipeline Stage 2 — Transformation",
            "Raw Layer  →  Curated Layer  (canonical, quality-checked, PII masked)")

    C1X, C1W = CX, Inches(4.75)
    C2X, C2W = CX+Inches(5.2), Inches(4.75)
    _vline(s, CX+Inches(5.0), BY, Inches(3.5), color=RULE_C)

    _blt(s, ["Reads Parquet from raw layer — streaming, no full load",
             "Field mapping: source → canonical names (v1.json)",
             "Transforms: CAST · CONCAT · DATE_FORMAT · DIRECT",
             "Missing fields: raise_error / use_default / drop_field",
             "Quality check: completeness, null rate, type conformance",
             "PII masking applied before curated write — never skipped",
             "SCD Type 1 merge: partition = full current state",
             "Lineage record emitted for every run",],
         C1X, BY, C1W, Inches(3.5),
         sz=Pt(12.5), color=TEXT, hdg="Processing Steps", hc=DEEP_BLUE)

    _blt(s, ["Mappings in S3 (curated bucket) — version-controlled",
             "Path: config/field_mappings/{source}/{entity}/v1.json",
             "No code change to add a new mapping — upload + run",
             "primary_key_field → upsert merge; absent = append",
             "soft_delete_field → tombstone (is_deleted flag kept)",
             "Quality thresholds configurable per entity",
             "Output: curated_date=YYYY-MM-DD/run_id=.../data.parquet",],
         C2X, BY, C2W, Inches(3.5),
         sz=Pt(12.5), color=TEXT, hdg="Configuration (no code needed)", hc=ORANGE)

    _box(s, CX, Inches(5.5), Inches(10.1), Inches(0.56), fill=RGBColor(0xF0,0xF4,0xF8))
    _t(s, '{ "source_fields":["AnnualRevenue"], "target_field":"annual_revenue", '
          '"transformation":"CAST","type":"decimal","missing_field_behavior":"use_default","default_value":"0" }',
       CX+Inches(0.1), Inches(5.56), Inches(9.9), Inches(0.48),
       sz=Pt(10), color=DEEP_BLUE, font=FM)

    _footer(s)


def s08_entity_res(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "Pipeline Stage 3 — Entity Resolution",
            "Cross-source deduplication  →  one Golden Record per real-world entity")

    for i, (c, nm, flds) in enumerate([
        (DEEP_BLUE, "Salesforce Account",
         "account_id\naccount_name\nbilling_country\nis_active · annual_revenue"),
        (SKY, "Sage Intacct Customer",
         "customer_id\ncompany_name\ncredit_limit\noutstanding_balance"),
        (GREEN, "Sage X3 Customer",
         "customer_code\ncustomer_name\ncountry_code\naccount_manager"),
    ]):
        sx = CX + i * Inches(3.4)
        bw = Inches(3.15)
        _box(s, sx, BY, bw, Inches(2.05), fill=c)
        _t(s, nm, sx+Inches(0.12), BY+Inches(0.1), bw-Inches(0.2), Inches(0.4),
           sz=Pt(13), bold=True, color=WHITE, font=FH)
        _t(s, flds, sx+Inches(0.12), BY+Inches(0.55), bw-Inches(0.2), Inches(1.4),
           sz=Pt(12), color=RGBColor(0xDD,0xEE,0xFF), font=FB)

    _t(s, "⬇  MATCH & MERGE  ⬇",
       CX+Inches(3.5), BY+Inches(2.12), Inches(3.0), Inches(0.42),
       sz=Pt(14), bold=True, color=ORANGE, align=PP_ALIGN.CENTER, font=FH)

    _box(s, CX, BY+Inches(2.62), Inches(10.1), Inches(1.45), fill=RGBColor(0xFF,0xF8,0xF0), lc=ORANGE, lw=Pt(1.5))
    _t(s, "🏆  GOLDEN RECORD  —  company entity",
       CX+Inches(0.15), BY+Inches(2.70), Inches(9.8), Inches(0.38),
       sz=Pt(14), bold=True, color=ORANGE, font=FH)
    _t(s, "entity_id  ·  account_name  ·  annual_revenue (Salesforce preferred)  ·  "
          "credit_limit (Sage preferred)  ·  outstanding_balance  ·  is_active\n"
          "Survivorship policy: configurable field-level source preference per entity type",
       CX+Inches(0.15), BY+Inches(3.1), Inches(9.8), Inches(0.85),
       sz=Pt(12), color=TEXT, font=FB)

    _blt(s, ["Blocking + fuzzy name matching — configurable per entity type",
             "Survivorship: preferred source per field (Salesforce → is_active; Sage → credit_limit)",
             "Entity types: company  ·  person  ·  contract  —  extensible via config",
             "Output → s3://edl-analytics-087972550871/canonical/{entity_type}/golden_date={date}/"],
         CX, BY+Inches(4.15), Inches(10.1), Inches(1.2),
         sz=Pt(12), color=TEXT)

    _footer(s)


def s09_analytics(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "Pipeline Stage 4 — Analytics Publisher",
            "Golden Records  →  Athena-queryable partitioned tables")

    _blt(s, ["Reads golden records from analytics layer (canonical/ prefix)",
             "Writes partitioned Parquet: analytics_date=YYYY-MM-DD",
             "Registers AWS Glue Data Catalog table",
             "Adds Hive partition for the run's analytics_date",
             "Grants Lake Formation table permissions to users",
             "CloudWatch metrics: record count, latency",],
         CX, BY, Inches(4.75), Inches(3.5),
         sz=Pt(13), color=TEXT, hdg="What it does", hc=DEEP_BLUE)

    _box(s, CX+Inches(5.1), BY, Inches(4.85), Inches(4.8), fill=RGBColor(0xF5,0xF8,0xFF))
    _t(s, "Query in Amazon Athena", CX+Inches(5.25), BY+Inches(0.1), Inches(4.55), Inches(0.35),
       sz=Pt(13), bold=True, color=DEEP_BLUE, font=FH)
    queries = [
        ("-- All companies", MUTED),
        ("SELECT * FROM edl_analytics.company", TEXT),
        ("  WHERE analytics_date='2026-07-02';", TEXT),
        ("", TEXT),
        ("-- Active contracts only", MUTED),
        ("SELECT COUNT(*) FROM edl_analytics.contract", TEXT),
        ("  WHERE analytics_date='2026-07-02'", TEXT),
        ("  AND is_deleted = false;", TEXT),
        ("", TEXT),
        ("-- Cross-entity join", MUTED),
        ("SELECT c.account_name, COUNT(k.contract_id)", TEXT),
        ("FROM edl_analytics.company c", TEXT),
        ("JOIN edl_analytics.contract k", TEXT),
        ("  ON c.account_id = k.tenant_id", TEXT),
        ("GROUP BY c.account_name;", TEXT),
    ]
    for i, (ln, c) in enumerate(queries):
        _t(s, ln, CX+Inches(5.25), BY+Inches(0.55)+i*Inches(0.22),
           Inches(4.55), Inches(0.22), sz=Pt(10), color=c, font=FM)

    _box(s, CX, BY+Inches(5.0), Inches(10.1), Inches(0.46), fill=DEEP_BLUE)
    _t(s, "Live today (dev):  34 companies  ·  49 persons  ·  35,971 contracts  — queryable via Athena right now",
       CX+Inches(0.15), BY+Inches(5.07), Inches(9.8), Inches(0.35),
       sz=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FH)

    _footer(s)


def s10_governance(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "Data Governance & Security",
            "Compliance built into the pipeline — not bolted on afterwards")

    panels = [
        (DEEP_BLUE, "PII Masking",
         "Applied at transformation — before curated write\n"
         "Sensitive fields tokenised, never plain text\n"
         "Policy version-controlled per entity\n"
         "Non-optional — pipeline fails if policy missing"),
        (SKY, "Data Lineage",
         "Emitted for every run:\nsource → curated → golden record\n"
         "Field count, run_id, entity_id, timestamp\n"
         "Queryable via Athena or S3"),
        (PURPLE, "Data Classification",
         "Fields: PUBLIC / INTERNAL / CONFIDENTIAL / PII\n"
         "Applied at source onboarding\n"
         "Drives masking + retention rules\n"
         "Registered in AWS Glue Catalog"),
        (ORANGE, "Retention Policy",
         "Configurable per classification tier\n"
         "Raw: S3 Object Lock (immutable)\n"
         "Curated/analytics: policy-enforced\n"
         "Dev: 30-day; Prod: 365-day logs"),
    ]
    pw = Inches(2.4)
    for (c, title, body), x in zip(panels, [CX, CX+Inches(2.55), CX+Inches(5.1), CX+Inches(7.65)]):
        _box(s, x, BY, pw, Inches(3.1), fill=c)
        _t(s, title, x+Inches(0.12), BY+Inches(0.12), pw-Inches(0.2), Inches(0.42),
           sz=Pt(13), bold=True, color=WHITE, font=FH)
        _box(s, x+Inches(0.08), BY+Inches(0.58), pw-Inches(0.16), Pt(1), fill=WHITE)
        _t(s, body, x+Inches(0.12), BY+Inches(0.72), pw-Inches(0.2), Inches(2.3),
           sz=Pt(11.5), color=RGBColor(0xEE,0xEE,0xEE), font=FB)

    _blt(s, ["All credentials via AWS Secrets Manager — never in code or logs",
             "IAM least-privilege: roles scoped to exact S3 prefixes and DynamoDB tables",
             "Only designated Lambda role can write to each S3 layer",
             "All event inputs validated with regex (OWASP A03 injection prevention)",
             "No record field values ever logged — counts/IDs only (OWASP A09)",],
         CX, BY+Inches(3.28), Inches(10.1), Inches(1.7),
         sz=Pt(12.5), color=TEXT, hdg="Security Controls", hc=CHARCOAL)

    _footer(s)


def s11_observability(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "Observability & Reliability",
            "Every run is measurable, traceable, and self-healing")

    obs = [("Structured Logging",     "JSON per stage: run_id, entity_id, counts — no raw field values"),
           ("CloudWatch Metrics",     "Custom namespace: record_count, latency_ms, quality_score"),
           ("Run Audit Log",          "Immutable DynamoDB record: start/end time, count, outcome, error"),
           ("Schema Drift Alerts",    "Drift report written to S3; CloudWatch alarm on breaking change"),
           ("Quality Report",         "Completeness + null rate per field after every transformation")]

    rel = [("Step Functions Retries", "Automatic retry with exponential backoff + jitter"),
           ("Dead-Letter Queue",      "Unrecoverable failures → SQS DLQ for human review and replay"),
           ("Circuit Breaker",        "DynamoDB-backed — stops hammering broken sources"),
           ("State Machine Guards",   "IsNull checks prevent 0-record extractions cascading"),
           ("Idempotency",            "Re-running same pipeline = identical output — safe to replay")]

    C1X, C1W = CX, Inches(4.75)
    C2X, C2W = CX+Inches(5.2), Inches(4.75)
    _vline(s, CX+Inches(5.0), BY, Inches(4.8), color=RULE_C)

    _t(s, "Monitoring", C1X, BY, C1W, Inches(0.35), sz=Pt(15), bold=True, color=DEEP_BLUE, font=FH)
    for i, (title, body) in enumerate(obs):
        ty = BY + Inches(0.42) + i * Inches(0.88)
        _box(s, C1X, ty, Inches(0.22), Inches(0.65), fill=DEEP_BLUE)
        _t(s, title, C1X+Inches(0.3), ty+Inches(0.02), C1W-Inches(0.32), Inches(0.28),
           sz=Pt(12.5), bold=True, color=DEEP_BLUE, font=FH)
        _t(s, body, C1X+Inches(0.3), ty+Inches(0.32), C1W-Inches(0.32), Inches(0.5),
           sz=Pt(12), color=TEXT, font=FB)

    _t(s, "Reliability", C2X, BY, C2W, Inches(0.35), sz=Pt(15), bold=True, color=ORANGE, font=FH)
    for i, (title, body) in enumerate(rel):
        ty = BY + Inches(0.42) + i * Inches(0.88)
        _box(s, C2X, ty, Inches(0.22), Inches(0.65), fill=ORANGE)
        _t(s, title, C2X+Inches(0.3), ty+Inches(0.02), C2W-Inches(0.32), Inches(0.28),
           sz=Pt(12.5), bold=True, color=ORANGE, font=FH)
        _t(s, body, C2X+Inches(0.3), ty+Inches(0.32), C2W-Inches(0.32), Inches(0.5),
           sz=Pt(12), color=TEXT, font=FB)

    _footer(s)


def s12_infra(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "Infrastructure & DevOps",
            "Terraform-managed  ·  Environment-parity  ·  One Lambda zip for all four functions")

    C1X, C1W = CX, Inches(4.75)
    C2X, C2W = CX+Inches(5.2), Inches(4.75)
    _vline(s, CX+Inches(5.0), BY, Inches(4.8), color=RULE_C)

    _blt(s, ["Terraform ≥ 1.8 — full infrastructure as code",
             "Modules: iam · storage · networking · lambda · orchestration",
             "Remote state: S3 + DynamoDB lock (single-writer safety)",
             "3 environments: dev / staging / prod — same config, different tfvars",
             "Env creation: init → iam → lambdas → orchestration",
             "Single Lambda zip serves all 4 functions (different handlers)",
             "Build: make lambda-package  →  Upload: make lambda-upload",],
         C1X, BY, C1W, Inches(4.5),
         sz=Pt(12.5), color=TEXT, hdg="Infrastructure", hc=DEEP_BLUE)

    _blt(s, ["Python 3.14, hatchling build system, pyproject.toml",
             "Tests: pytest + moto (AWS mocks — no real AWS calls)",
             "Lint: ruff  ·  Types: mypy  ·  Security: bandit",
             "Dependency CVE scan: pip-audit",
             "Coverage gate: 80% — enforced in CI",
             "Local dry-run scripts — test connectivity, no S3 write",
             "CI: lint → type-check → test → security (one command)",],
         C2X, BY, C2W, Inches(4.5),
         sz=Pt(12.5), color=TEXT, hdg="Developer Experience & CI", hc=ORANGE)

    _footer(s)


def s13_live(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "Live Data — Dev Environment",
            "Real business data flowing end-to-end  ·  as of July 2026")

    metrics = [(GREEN_OK,  "35,971+",  "Contracts\n(MySQL RDS)"),
               (DEEP_BLUE, "34",       "Companies\n(Salesforce)"),
               (SKY,       "49",       "Persons\n(Salesforce)"),
               (PURPLE,    "4",        "Sources\nconnected"),
               (ORANGE,    "< 4 hrs",  "End-to-end\nlatency"),
               (GREEN,     "0",        "Mapping\nfailures")]
    mw = Inches(10.1) / len(metrics)
    for i, (c, num, lbl) in enumerate(metrics):
        mx = CX + i * mw
        _box(s, mx, BY, mw-Inches(0.1), Inches(1.95), fill=c)
        _t(s, num, mx, BY+Inches(0.12), mw-Inches(0.1), Inches(0.95),
           sz=Pt(32), bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FH)
        _t(s, lbl, mx, BY+Inches(1.1), mw-Inches(0.1), Inches(0.75),
           sz=Pt(11), color=RGBColor(0xEE,0xEE,0xEE), align=PP_ALIGN.CENTER, font=FB)

    run_cx = [CX, CX+Inches(2.8), CX+Inches(4.7), CX+Inches(6.35), CX+Inches(9.3)]
    run_cw = [Inches(2.65), Inches(1.75), Inches(1.5), Inches(2.8), Inches(0.68)]
    run_hdrs = ["Entity", "Load type", "Records", "Athena table", "Status"]
    ht = BY + Inches(2.12)
    for h, x, w in zip(run_hdrs, run_cx, run_cw):
        _box(s, x, ht, w, Inches(0.36), fill=DEEP_BLUE)
        _t(s, h, x+Inches(0.07), ht+Inches(0.06), w-Inches(0.1), Inches(0.28),
           sz=Pt(11), bold=True, color=WHITE, font=FH)

    runs = [("MySQL — Contracts",     "Incremental (SCD1)", "35,971",  "edl_analytics.contract",        "✅"),
            ("Salesforce Account",    "Full load",          "34",       "edl_analytics.company",          "✅"),
            ("Salesforce Contact",    "Incremental",        "49",       "edl_analytics.person",           "✅"),
            ("Sage Intacct Customer", "Incremental",        "Live",     "edl_analytics.company (merged)", "✅"),
            ("Sage Intacct Vendor",   "Incremental",        "Live",     "Curated layer",                      "✅"),]
    for ri, (rv, bg) in enumerate(zip(runs, [WHITE, LIGHT_BG]*10)):
        rt = BY + Inches(2.48) + ri * Inches(0.44)
        for val, x, w in zip(rv, run_cx, run_cw):
            _box(s, x, rt, w, Inches(0.42), fill=bg)
            c = GREEN_OK if val == "✅" else TEXT
            _t(s, val, x+Inches(0.07), rt+Inches(0.06), w-Inches(0.1), Inches(0.35),
               sz=Pt(11), color=c, font=FB,
               align=PP_ALIGN.CENTER if val == "✅" else PP_ALIGN.LEFT)

    _footer(s)


def s14_roi(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "Business Outcomes & ROI",
            "Measurable improvements delivered — compared to before the platform")

    cx  = [CX, CX+Inches(3.3), CX+Inches(6.35)]
    cw  = [Inches(3.3), Inches(3.05), Inches(3.75)]
    for h, x, w, c in zip(["Metric","Before","After"], cx, cw, [CHARCOAL, RED, GREEN_OK]):
        _box(s, x, BY, w, Inches(0.4), fill=c)
        _t(s, h, x+Inches(0.1), BY+Inches(0.08), w-Inches(0.15), Inches(0.3),
           sz=Pt(12), bold=True, color=WHITE, font=FH)

    rows = [
        ("Time to data availability", "24–72 hours (manual)",              "1–4 hours (automated)"),
        ("Customer identity",         "3 disconnected views per system",   "Single golden record"),
        ("PII in analytics",          "Uncontrolled",                       "Masked at pipeline level"),
        ("Audit trail",               "None",                               "Full lineage: source→golden"),
        ("New source onboarding",     "2–4 weeks (code + deploy)",          "2–3 days (config only)"),
        ("Credential security",       "Scripts + shared .env files",        "AWS Secrets Manager"),
        ("Compliance",                "Manual documentation",               "Automated lineage + retention"),
        ("Data quality",              "No monitoring",                      "Report per entity per run"),
    ]
    for ri, (rv, bg) in enumerate(zip(rows, [WHITE, LIGHT_BG]*10)):
        rt = BY + Inches(0.4) + ri * Inches(0.5)
        for val, x, w in zip(rv, cx, cw):
            _box(s, x, rt, w, Inches(0.48), fill=bg)
            _t(s, val, x+Inches(0.1), rt+Inches(0.07), w-Inches(0.15), Inches(0.38),
               sz=Pt(11.5), color=TEXT, font=FB)

    _box(s, CX, BY+Inches(4.5), Inches(10.1), Inches(0.82), fill=RGBColor(0xF0,0xF4,0xF8))
    _t(s, "AWS Infrastructure: ~$469/month",
       CX+Inches(0.15), BY+Inches(4.58), Inches(3.2), Inches(0.32),
       sz=Pt(13), bold=True, color=DEEP_BLUE, font=FH)
    _t(s, "vs. Fivetran SaaS: $3,000–$5,000/month per source  |  Payback: < 1 month  |  Plus 40–60 hrs/month avoided labour",
       CX+Inches(3.5), BY+Inches(4.58), Inches(6.5), Inches(0.32),
       sz=Pt(12), color=TEXT, font=FB)
    _t(s, "10× cost saving vs. SaaS — with full ownership",
       CX+Inches(0.15), BY+Inches(4.92), Inches(9.5), Inches(0.3),
       sz=Pt(12), bold=True, color=ORANGE, font=FH)

    _footer(s)


def s15_stack(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "Technology Stack",
            "Purpose-built on proven AWS-native services and open standards")

    cats = [
        (DEEP_BLUE, "Runtime",
         ["Python 3.14  —  AWS Lambda", "PyArrow  —  Parquet I/O",
          "Pydantic  —  config validation", "boto3  —  AWS SDK"]),
        (SKY, "Orchestration",
         ["AWS Step Functions", "Amazon EventBridge Scheduler",
          "AWS Lambda (4 functions, 1 zip)"]),
        (PURPLE, "Storage",
         ["Amazon S3 (3 data layers)", "Amazon DynamoDB (3 tables)",
          "AWS Glue Data Catalog"]),
        (RGBColor(0xC0,0x55,0x10), "Analytics",
         ["Amazon Athena (SQL queries)", "AWS Lake Formation (access)",
          "Apache Parquet (columnar)"]),
        (RED, "Security & Ops",
         ["AWS Secrets Manager", "AWS IAM (least-privilege)",
          "Amazon CloudWatch", "Amazon SQS (dead-letter)"]),
        (GREEN, "DevOps",
         ["Terraform ≥ 1.8 (IaC)", "pytest + moto (testing)",
          "ruff + mypy + bandit", "GitHub (source control)"]),
    ]
    cw = Inches(1.62)
    for (c, cat, items), x in zip(cats, [CX+i*Inches(1.7) for i in range(6)]):
        _box(s, x, BY, cw, Inches(5.1), fill=c)
        _t(s, cat, x+Inches(0.1), BY+Inches(0.1), cw-Inches(0.16), Inches(0.42),
           sz=Pt(12), bold=True, color=WHITE, font=FH)
        _box(s, x+Inches(0.08), BY+Inches(0.56), cw-Inches(0.16), Pt(1), fill=WHITE)
        for i, item in enumerate(items):
            _t(s, f"• {item}", x+Inches(0.1), BY+Inches(0.68)+i*Inches(0.55),
               cw-Inches(0.16), Inches(0.52), sz=Pt(11),
               color=RGBColor(0xEE,0xEE,0xEE), font=FB)

    _footer(s)


def s16_roadmap(prs):
    s = _blank(prs); _white(s); _strip(s)
    _header(s, "What's Next — Roadmap",
            "Dev is live  ·  Staging and Production are the immediate priorities")

    lanes = [
        (GREEN_OK, "Now — Dev ✅",
         ["All 4 Lambda stages live & tested",
          "Salesforce, MySQL, Sage connected",
          "35,971+ records in analytics layer",
          "Athena queries operational",
          "SCD Type 1 incremental merge",
          "Entity resolution: golden records live"]),
        (AMBER, "Next — Staging",
         ["Provision staging DynamoDB tables",
          "Terraform: iam → lambdas → orchestration",
          "Upload Lambda zip to staging bucket",
          "Seed entity configs to staging",
          "End-to-end smoke test",
          "Staging sign-off before prod"]),
        (DEEP_BLUE, "Soon — Production",
         ["Prod infrastructure provisioning",
          "365-day log retention (in HCL already)",
          "Credential rotation policy enabled",
          "Monitoring dashboards + alerts",
          "Business data quality sign-off",
          "Go-live readiness checklist"]),
        (ORANGE, "Backlog",
         ["NetSuite ERP connector",
          "Sage X3 entity resolution",
          "Sage Intacct AR Invoice / AP Bill",
          "Lambda timeout circuit breaker",
          "value_map transform for booleans",
          "Serving store loader (write-back)"]),
    ]
    lw = Inches(2.35)
    for (c, title, items), x in zip(lanes, [CX+i*Inches(2.5) for i in range(4)]):
        _box(s, x, BY, lw, Inches(0.42), fill=c)
        _t(s, title, x+Inches(0.1), BY+Inches(0.06), lw-Inches(0.16), Inches(0.35),
           sz=Pt(12), bold=True, color=WHITE, font=FH)
        _box(s, x, BY+Inches(0.42), lw, Inches(4.8), fill=LIGHT_BG, lc=c, lw=Pt(1))
        for i, item in enumerate(items):
            _t(s, f"▸  {item}", x+Inches(0.12), BY+Inches(0.55)+i*Inches(0.76),
               lw-Inches(0.2), Inches(0.72), sz=Pt(12), color=TEXT, font=FB)

    _footer(s)


def s17_closing(prs):
    s = _blank(prs); _white(s); _strip(s)

    _t(s, "A Platform Built\nfor the Long Term",
       CX, Inches(0.42), Inches(9.5), Inches(1.8),
       sz=Pt(42), bold=True, color=CHARCOAL, font=FH)

    _rule(s, Inches(2.3), color=DEEP_BLUE, thick=Pt(3))

    items = [
        (DEEP_BLUE, "Automated, end-to-end pipeline — running in the dev environment today"),
        (ORANGE,    "Four connected sources — extensible to any REST or SQL source via config"),
        (GREEN_OK,  "Entity resolution delivering one golden record per customer across all systems"),
        (SKY,       "Governed by design: lineage, PII masking, classification, retention built in"),
        (DEEP_BLUE, "Queryable in standard SQL via Amazon Athena — no exports, no waiting"),
        (ORANGE,    "Infrastructure as code — reproducible across dev, staging, and production"),
        (GREEN_OK,  "10× lower cost than SaaS alternatives — with full ownership and control"),
    ]
    for i, (c, text) in enumerate(items):
        ty = Inches(2.45) + i * Inches(0.62)
        _box(s, CX, ty, Inches(0.32), Inches(0.46), fill=c)
        _t(s, "✅", CX, ty+Inches(0.04), Inches(0.32), Inches(0.4),
           sz=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)
        _t(s, text, CX+Inches(0.42), ty+Inches(0.05), Inches(9.5), Inches(0.46),
           sz=Pt(15), color=TEXT, font=FB)

    _t(s, "Enterprise Data Lake Platform  ·  Platform Engineering  ·  July 2026",
       CX, FTY, Inches(9.0), Inches(0.32), sz=Pt(11), color=MUTED, font=FB)


def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    s01_title(prs)
    s02_problem(prs)
    s03_overview(prs)
    s04_architecture(prs)
    s05_sources(prs)
    s06_extraction(prs)
    s07_transformation(prs)
    s08_entity_res(prs)
    s09_analytics(prs)
    s10_governance(prs)
    s11_observability(prs)
    s12_infra(prs)
    s13_live(prs)
    s14_roi(prs)
    s15_stack(prs)
    s16_roadmap(prs)
    s17_closing(prs)

    out = os.path.join(os.path.dirname(os.path.dirname(__file__)),
                       "Enterprise_Data_Lake_Platform.pptx")
    prs.save(out)
    print(f"✅  Saved: {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
